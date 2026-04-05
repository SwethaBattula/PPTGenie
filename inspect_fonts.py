from pptx import Presentation

prs = Presentation("PPTGenie_Output.pptx")

with open("output_fonts.txt", "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"\n--- Slide {i+1} ---\n")
        f.write(f"Title: {slide.shapes.title.text}\n")
        try:
            content = slide.placeholders[1]
            for j, p in enumerate(content.text_frame.paragraphs):
                runs = p.runs
                if runs:
                    f.write(f"  Para {j}: font size = {runs[0].font.size}, text = {p.text[:20]}...\n")
                else:
                    f.write(f"  Para {j}: NO RUNS (empty)\n")
        except Exception as e:
            f.write(f"  Error reading content: {e}\n")
