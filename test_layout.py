from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Test Title"

content = slide.placeholders[1]
content.left = Inches(0.5)
content.top = Inches(1.5)
content.width = Inches(4.5)
content.height = Inches(5.0)
content.text = "Testing bullet points\nPoint 2"

prs.save("test_layout.pptx")
print("Success")
