# Importing Presentation class from python-pptx library
# This library is used to create and manipulate PowerPoint files programmatically
from pptx import Presentation
import datetime
import requests
import os
from pptx.util import Pt, Inches


class PPTTools:
    """
    PPTTools acts as the MCP (Model Context Protocol) tool layer.

    This class provides functions (tools) that the agent can call to perform actions
    such as creating a presentation, adding slides, and writing content.

    The agent DOES NOT directly manipulate the PowerPoint file.
    Instead, it interacts with these tools, maintaining separation between:
    - Decision-making (Agent)
    - Execution (Tools)
    """


    def __init__(self):
        """
        Initializes the PPTTools object.

        Why:
        The presentation object is initialized as None and will be created
        only when the agent explicitly calls the create_presentation tool.
        This ensures controlled execution flow.
        """
        self.prs = None


    def create_presentation(self):
        """
        Tool: create_presentation

        Purpose:
        Initializes a new PowerPoint presentation.

        Why:
        The agent must explicitly start a presentation before adding slides.
        This prevents invalid operations (like adding slides without a file).

        Returns:
        A confirmation message indicating successful creation.
        """
        self.prs = Presentation()
        return "Presentation created successfully"
    
    
    def add_slide(self, title):
        """
        Tool: add_slide

        Purpose:
        Adds a new slide to the presentation with a given title.

        Why:
        The agent uses this tool to dynamically create slides based on its plan.
        This ensures slides are not hardcoded and are generated at runtime.

        Parameters:
        title (str): The title of the slide.

        Returns:
        The created slide object so that further content can be added.
        """

        # Check if presentation exists before adding slides
        # This prevents runtime errors and enforces correct tool usage order
        if self.prs is None:
            raise Exception("Presentation not created. Call create_presentation first.")

        # Using layout index 1 (Title + Content layout)
        layout = self.prs.slide_layouts[1]

        # Adding a new slide using the selected layout
        slide = self.prs.slides.add_slide(layout)

        # Setting the title of the slide
        title_shape = slide.shapes.title
        title_shape.text = title
        
        from pptx.util import Pt
        
        # Explicitly setting a fixed font size for the title across all slides
        # This prevents PowerPoint's AutoFit from making short titles huge and long titles tiny
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(36)
        self.current_slide = slide   # ⭐ store latest slide
        return slide
    
    def generate_content(self, title):
     """
     Tool: generate_content

     Purpose:
     Uses Gemini (LLM) to generate meaningful bullet points.

     Why:
     Replaces generic hardcoded text with intelligent content.
     """

     from google import genai

     client = genai.Client(api_key="YOUR_API_KEY")

     prompt = f"""
       Generate 4 high-quality bullet points for a PowerPoint slide on: {title}

        Rules:
        - Be specific and informative
        - Avoid generic phrases like "it explains"
        - Each point should be short and clear
        - No numbering
     """

     response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt
     )

     #Convert response → list
     points = response.text.strip().split("\n")

     return points


    def write_content(self, bullet_points, has_image=False):
        """
        Tool: write_content

        Purpose:
        Writes bullet point content into a given slide.

        Why:
        The agent generates textual content for each slide dynamically.
        This tool ensures that the content is inserted in a structured
        and consistent format.

        Parameters:
        slide: The slide object returned from add_slide.
        bullet_points (list): A list of strings representing bullet points.
        has_image (bool): Determines if text should be constrained to left area.

        Returns:
        A confirmation message indicating content has been added.
        """
        # ✅ Get the latest slide created
        slide = self.current_slide

        # Access content placeholder
        text_frame = slide.placeholders[1].text_frame

        # Clear default text
        text_frame.clear()

        content_placeholder = slide.placeholders[1]
        content_placeholder.left = Inches(0.5)
        content_placeholder.top = Inches(2.2)

        if has_image:
         content_placeholder.width = Inches(4.5)
        else:
         content_placeholder.width = Inches(9.0)

        content_placeholder.height = Inches(4.3)

        bullet_points = bullet_points[:4]

        for i, point in enumerate(bullet_points):
            if i == 0:
             paragraph = text_frame.paragraphs[0]
            else:
             paragraph = text_frame.add_paragraph()

            paragraph.text = point

        for run in paragraph.runs:
             run.font.size = Pt(18)

        return "Content added successfully"
    
    
    def generate_image(self, title):
     url = f"https://source.unsplash.com/800x600/?{title}"
     file_name = f"{title.replace(' ', '_')}.jpg"

     response = requests.get(url)

     with open(file_name, "wb") as f:
        f.write(response.content)

     return os.path.abspath(file_name)
    
    
    def generate_summary(self, title, bullet_points):
        """
        Tool: generate_summary

        Purpose:
        Generates a one-line summary for a slide based on its title and content.

        Why:
        A concise summary helps reinforce the key message of the slide,
        making the presentation more impactful and easier to understand.

        This also demonstrates enhanced agent capability beyond basic content generation.

        Parameters:
        title (str): The title of the slide.
        bullet_points (list): The bullet points of the slide.

        Returns:
        str: A single-line summary representing the slide content.
        """

        # Creating a simple summary by combining title and main idea
        # This can later be replaced with an LLM-based summarization
        if bullet_points:
            # Just take the first point to keep it one-line, avoiding repetitive titles
            summary = bullet_points[0]
        else:
            summary = f"Key Takeaway: {title} overview"

        return summary
    
    def add_image_to_slide(self, slide, image_path):
        """
        Tool: add_image_to_slide

        Purpose:
        Adds an image to a given slide.

        Why:
        Visual elements enhance the presentation quality and improve understanding.
        This tool allows the agent to insert images dynamically based on content.

        Parameters:
        slide: The slide object where the image will be added.
        image_path (str): Path to the image file.

        Returns:
        A confirmation message indicating the image has been added.
        """

        # Importing Inches for positioning the image on the slide
        from pptx.util import Inches

        # Positioning the image (left, top, width, height)
        # These values ensure the image does not overlap with text on the left
        left = Inches(5.25)   # strictly occupy right side
        top = Inches(2.2)     # lowered to avoid title
        width = Inches(4.25)  # strictly right-side width
        height = Inches(4.3)  # aligned with text block

        # Adding image to the slide
        slide.shapes.add_picture(image_path, left, top, width, height)

        return "Image added successfully"
    
    def save_presentation(self, filename = f"PPTGenie_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"):
        """
        Tool: save_presentation

        Purpose:
        Saves the PowerPoint presentation to the local disk.

        Why:
        This is the final step in the workflow where all created slides,
        content, and images are compiled into a single .pptx file.

        The agent calls this tool after completing all slide generation tasks.

        Parameters:
        filename (str): Name of the output file.

        Returns:
        str: Confirmation message with file name.
        """

        # Ensuring that a presentation exists before saving
        if self.prs is None:
            raise Exception("No presentation found. Create one before saving.")

        # Saving the presentation to disk
        self.prs.save(filename)

        return f"Presentation saved as {filename}"

    def add_summary_to_footer(self, slide, summary_text):
        """
        Tool: add_summary_to_footer

        Purpose:
        Adds a single-line summary string as a centered footer at the bottom of the slide.

        Why:
        Creates a structurally clean and distinct area for the slide summary,
        so it doesn't mix with bullet points.

        Parameters:
        slide: The slide object where the summary will be added.
        summary_text (str): The summary text to write.

        Returns:
        A confirmation message indicating summary has been added.
        """
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        # Positioning at the bottom center
        left = Inches(0.5)
        top = Inches(6.8)
        width = Inches(9.0)
        height = Inches(0.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = summary_text
        p.alignment = PP_ALIGN.CENTER
        
        for run in p.runs:
            run.font.size = Pt(12)  # Reduced to ensure one-line fit
            
        return "Summary added successfully"